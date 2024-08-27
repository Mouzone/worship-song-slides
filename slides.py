from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
import lyrics

# search for song, then get order and dictionary
# create slide
# -- some kind of gui to modify text content
# ---- modify font size
# ---- create slides and delete slides as well

song_name = input("Enter song name: ")
artist = input("Enter artist name: ")
order, lyrics_by_section = lyrics.search(song_name, artist)

ppt = Presentation()
blank_slide_layout = ppt.slide_layouts[6]
width = height = Inches(10)
left = (ppt.slide_width - width) / 2
top = (ppt.slide_height - height) / 2
for section_name in order:
    slide = ppt.slides.add_slide(blank_slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame

    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    run = p.add_run()
    run.text = lyrics_by_section[section_name]
    font = run.font
    font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


ppt.save(f"{song_name}.pptx")

